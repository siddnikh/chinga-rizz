import os
import json
from langchain.prompts.chat import (
    SystemMessagePromptTemplate,
    ChatPromptTemplate
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

async def generateIndex(q: str, slide_no: int):
    prompt_array = get_prompt_array("indexGeneration")
    prompt = ChatPromptTemplate(prompt_array)

    #TODO
    # llm = None
    # chain = LLMChain(llm=llm, prompt=prompt)
    # result = chain.run(user_input=f"Topic: {q}, Number of slides: {slide_no}")
    result = "yes"
    return result

async def generateDocContent(q: str, slide_no: int, index: list[str]):
    prompt_array = get_prompt_array("documentGeneration")
    prompt = ChatPromptTemplate(prompt_array)

    #TODO
    # llm = None
    # chain = LLMChain(llm=llm, prompt=prompt)
    # result = chain.run(user_input=f"Topic: {q}, Number of slides: {slide_no}, Index: {index}")
    result = "yes"
    return result

async def generateDoc(index: list[str], content: dict, generated_images: list[str]):
    
    prs = Presentation()
    # Set default font properties
    default_font = prs.theme.fonts[0]
    default_font.name = "Montserrat"
    default_font.size = Pt(12) 
    default_font.color.rgb = RGBColor(0, 0, 0)

    blank_layout = prs.slide_layouts[6]
    
    # Slide 1
    slide_1 = prs.slides.add_slide(blank_layout)
    title_1 = slide_1.shapes.title
    title_1.text = index[0]
    title_1.font_size = Pt(24)

    body_text_1 = slide_1.shapes.add_textbox(Inches(0.1), Inches(2), Inches(8), Inches(4))
    tf = body_text_1.text_frame
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.text = content[0]["content"]

    image = slide_1.shapes.add_picture(generated_images[0], Inches(0), Inches(2), Inches(3.33), Inches(4))

    # Resize and crop image to fit one-third of the slide width while maintaining aspect ratio
    image.crop(0, 0, image.width / image.height, 1)
    image.left = Inches(0)
    image.width = Inches(3.33)


def get_prompt_template_array(prompts: list[str]):
    prompt_array = []

    for prompt in prompts:
        prompt_array.append(SystemMessagePromptTemplate.from_template(prompt))

    return prompt_array

# category can be either indexGeneration, or documentGeneration, or imageGeneration
def get_prompt_array(category: str):
    
    file_name = "prompts.json"
    file_path = os.path.join(os.getcwd(), file_name)
    f = open(file_path)
    data = json.load(f)
    f.close()

    return get_prompt_template_array(data[category])
